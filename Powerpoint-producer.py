#!/usr/bin/env python

'''
Powerpoint-producer.py
A. J. McCulloch, February 2020
'''

####################################################################################################
# Import modules
####################################################################################################

from pptx import Presentation
from pptx.util import Cm, Pt
from datetime import date

####################################################################################################
# Define functions
####################################################################################################

"""
The tricky bit of all of this is addressing the correct placeholder.
For the UTAS template, I have pulled the placeholders via the function find_ph.
Sample output is shown later
"""

def find_ph(file):
    prs = Presentation(file)
    for cnt, s in enumerate(prs.slide_layouts):
        slide = prs.slides.add_slide(s)
        print(s.name)
        for shape in slide.placeholders:
            print('%d %s' % (shape.placeholder_format.idx, shape.name))
        if cnt != len(prs.slide_layouts)-1:
            print('')

# Function to delete slides
def delete_slide(prs, slide):
    #Make dictionary with necessary information
    id_dict = { slide.id: [i, slide.rId] for i,slide in enumerate(prs.slides._sldIdLst) }
    slide_id = slide.slide_id
    prs.part.drop_rel(id_dict[slide_id][1])
    del prs.slides._sldIdLst[id_dict[slide_id][0]]

####################################################################################################
####################################################################################################
# Code starts here
####################################################################################################
####################################################################################################

file = 'UTAS_PPT_Template.pptx' # Template file
# find_ph(file) # Get the placeholders

"""
For the template 'UTAS_PPT_Template.pptx',
find_ph returned the following:

Title Page
0 Title 1
1 Subtitle 2
10 Text Placeholder 3

Section Title
0 Title 1

Deductive – storyline on a page
0 Title 2
17 Text Placeholder 9
28 Text Placeholder 3
29 Text Placeholder 4
30 Text Placeholder 5
51 Text Placeholder 6
52 Text Placeholder 1
53 Text Placeholder 7
54 Text Placeholder 8

1_Chart
0 Title 1
17 Text Placeholder 4
19 Text Placeholder 2
20 Text Placeholder 3
23 Text Placeholder 5
24 Picture Placeholder 6

Blank
0 Title 1
17 Text Placeholder 2

End Page - Light
0 Title 3
1 Subtitle 1
10 Text Placeholder 2
"""

prs = Presentation('UTAS_PPT_Template.pptx') # Load the template powerpoint

####################################################################################################
# Title slide
####################################################################################################

title_layout = prs.slide_layouts[0] # Set the slide layout
slide = prs.slides.add_slide(title_layout) # Add the slide
# Title slide attributes
title = slide.shapes.title # Title
subtitle = slide.placeholders[1] # Subtitle
date_created = slide.placeholders[10] # Date

title.text = "This is the title text"
subtitle.text = "This is the subtitle"
date_created.text = "{:%d/%m/%Y}".format(date.today())

####################################################################################################
# Section slide
####################################################################################################

section_layout = prs.slide_layouts[1] # Set the slide layout
slide = prs.slides.add_slide(section_layout) # Add the slide
# Section slide attribute
section = slide.shapes.title # Title

section.text = "This is the section title"

####################################################################################################
# Storyline slide
####################################################################################################

storyline_layout = prs.slide_layouts[2] # Set the slide layout
slide = prs.slides.add_slide(storyline_layout) # Add the slide
# Storyline slide attributes
storyline = slide.shapes.title # Title
question = slide.placeholders[51]
situation = slide.placeholders[54]
complication = slide.placeholders[53]
resolution = slide.placeholders[52]
situation_overview = slide.placeholders[30]
complication_overview = slide.placeholders[28]
resolution_overview = slide.placeholders[29]
source = slide.placeholders[17]

storyline.text = "This is the story"
question.text = "This is a question"
situation.text = "This is a situation"
complication.text = "This is a complication"
resolution.text = "This is a resolution"
situation_overview.text = "This is a situation overview"
complication_overview.text = "This is a complication overview"
resolution_overview.text = "This is a resolution overview"
source.text = "This is a source"

####################################################################################################
# Chart slide
####################################################################################################

chart_layout = prs.slide_layouts[3] # Set the slide layout
slide = prs.slides.add_slide(chart_layout) # Add the slide
# Chart slide attributes
title = slide.shapes.title # Title
chart_title = slide.placeholders[23]
picture = slide.placeholders[24]
takeaway = slide.placeholders[20]
overview = slide.placeholders[19]
source = slide.placeholders[17]

title.text = "This is a title"
chart_title.text = "This is a chart title"
#chart = slide.placeholders[18]
takeaway.text = "This is the takeaway text"
overview.text = "This is the overview text"
source.text = "This is the source"


img_path = 'Placeholder.jpg'
picture = picture.insert_picture(img_path)

#left = Cm(5)
#height = Cm(5.5)
#pic = slide.shapes.add_picture(img_path, left, top, height=height)

####################################################################################################
# Blank slide
####################################################################################################

chart_layout = prs.slide_layouts[4] # Set the slide layout
slide = prs.slides.add_slide(chart_layout) # Add the slide

title = slide.shapes.title # Title
source = slide.placeholders[17] # Source

title.text = "This is a title"
source.text = "This is the source"

####################################################################################################
# Ending slide
####################################################################################################

title_layout = prs.slide_layouts[0] # Set the slide layout
slide = prs.slides.add_slide(title_layout) # Add the slide
# Title slide attributes
title = slide.shapes.title # Title
subtitle = slide.placeholders[1] # Subtitle
date_created = slide.placeholders[10] # Date

title.text = "This is the title text"
subtitle.text = "This is the subtitle"
date_created.text = "{:%d/%m/%Y}".format(date.today())

####################################################################################################
# Output
####################################################################################################

# Delete the original slides with no content
for cnt, slide in enumerate(prs.slides):
    if cnt < 6:
        delete_slide(prs, slide)

# Save the output file
prs.save('test.pptx')
